from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fitz
import os
from collections import defaultdict
from ai.module.util import FileUtil

# Module-level constants
IMAGE_EXTENSIONS = ['png', 'jpeg', 'jpg']
IMAGE_SIZE_THRESHOLD = (200, 200)

class ExtractInfo:
    def __init__(self):
        self.fileutil = FileUtil()
        self.reset()

    def reset(self):
        self.file_name = None
        self.name = None
        self.file_ext = None
        self.text_dict = {'total': [], 'page': defaultdict(list)}
        self.img_dict = defaultdict(list)
        self.page_num = None
        self.img_num = 0
        self.image_blob = []

class PPTInfoExtract(ExtractInfo):
    def __init__(self):
        super().__init__()

    def reset(self):
        super().reset()
        self.parsed = None

    def run(self, file_name):
        self._extract(file_name)
        res = {}
        texts = self.text_dict['page']
        imgs = self.img_dict

        pages = set(list(texts.keys()) + list(imgs.keys()))
        for page in range(min(pages), max(pages) + 1):
            res[page] = {}
            res[page]['text'] = texts[page]
            res[page]['img'] = imgs[page]
        return res

    def _extract(self, file_name):
        self.reset()
        self.file_name = file_name
        self.name, self.file_ext = os.path.splitext(file_name)
        path = os.path.join(self.fileutil.orignal_dir, file_name)
        self.parsed = Presentation(path)

        for idx, slide in enumerate(self.parsed.slides, start=1):
            self.page_num = idx
            for shape in slide.shapes:
                self._info_extract(shape)

    def _info_extract(self, shape):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                self._info_extract(s)
        else:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_blob = shape.image.blob

                    if image_blob not in self.image_blob:
                        ext = shape.image.ext
                        size = shape.image.size
                        if ext in IMAGE_EXTENSIONS and size[0] >= IMAGE_SIZE_THRESHOLD[0] and size[1] >= IMAGE_SIZE_THRESHOLD[1]:
                            self.img_num += 1
                            num = str(self.img_num).zfill(3)
                            save_path = os.path.join(self.fileutil.res_dir, f"{self.name}/image_{num}.{ext}")
                            with open(save_path, "wb") as file:
                                file.write(image_blob)
                            self.image_blob.append(image_blob)
                            self.img_dict[self.img_num].append(save_path)
                except Exception:
                    pass
            else:
                if shape.has_text_frame:
                    if shape.text.strip() != "":
                        self.text_dict['page'][self.page_num].append(shape.text.strip())


class DOCXInfoExtract(ExtractInfo):
    def __init__(self) -> None:
        super().__init__()

    def reset(self):
        super().reset()
        self.doc = None

    def run(self, file_name):
        """_summary_

        Args:
            file_name (str): 문서 이름

        Returns:
            dict:
                page_num (int):
                    text (str): list
                    img (str): list
        """
        self._extract(file_name)
        res = {}
        texts = self.text_dict['page']
        imgs = self.img_dict

        self.text_dict['page'][self.page_num]
        self.text_dict['page'][self.page_num].append(
            self._extract_tables())
        self._extract_images()

        pages = set(list(texts.keys()) + list(imgs.keys()))
        for page in range(min(pages), max(pages)+1):
            res[page] = {}
            res[page]['text'] = texts[page]
            res[page]['img'] = imgs[page]
        return res

    def _extract(self, file_name):
        # path = "deep_learning_intro.pptx"
        self.reset()
        self.file_name = file_name
        self.name, self.file_ext = os.path.splitext(file_name)
        path = self.fileutil.orignal_dir + file_name
        self.doc = Document(path)

        self.text_dict['page'][self.page_num]
        self.text_dict['page'][self.page_num].append(
            self._extract_tables())
        self._extract_images()

    def _extract_text(self):
        paragraphs = [p.text for p in self.doc.paragraphs]
        return '\n'.join(paragraphs)

    def _extract_tables(self):
        tables = []
        for table in self.doc.tables:
            data = []
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    cells.append(cell.text)
                data.append(' '.join(cells))
            tables.append(' '.join(data))
        return '\n'.join(tables)

    def _extract_images(self):
        for rel in self.doc.part.rels.values():
            if "image" in rel.reltype:
                image = rel.target_part.image
                image_blob = image.blob

                # 이미지 중복 체크
                if image_blob not in self.image_blob:
                    ext = image.ext
                    size = image.px_height, image.px_width
                    if ext in ['png', 'jpeg', 'jpg'] and size[0] >= 200 and size[1] >= 200:
                        self.img_num += 1
                        num = str(self.img_num).zfill(3)
                        save_path = f"{self.fileutil.res_dir}{self.name}/image_{num}.{ext}"

                        with open(save_path, "wb") as file:
                            file.write(image_blob)
                        self.image_blob.append(image_blob)
                        self.img_dict[self.img_num].append(save_path)


# https://neurondai.medium.com/how-to-extract-text-from-a-pdf-using-pymupdf-and-python-caa8487cf9d
class PDFInfoExtract(ExtractInfo):
    def __init__(self):
        super().__init__()

    def reset(self):
        super().reset()
        self.doc = None

    def run(self, file_name):
        self._extract(file_name)
        res = {}
        texts = self.text_dict['page']
        imgs = self.img_dict

        pages = set(list(texts.keys()) + list(imgs.keys()))
        for page in range(min(pages), max(pages) + 1):
            res[page] = {}
            res[page]['text'] = texts[page]
            res[page]['img'] = imgs[page]
        return res

    def _extract(self, file_name):
        self.reset()
        self.file_name = file_name
        self.name, self.file_ext = os.path.splitext(file_name)
        path = os.path.join(self.fileutil.orignal_dir, file_name)
        self.pdf = fitz.open(path)
        # print(len(self.pdf))
        for idx, page_content in enumerate(self.pdf, start=1):
            self.page_num = idx
            # 텍스트 추출
            self._extract_text(page_content)
            # 이미지 추출
            self._extract_image(page_content)

    def _extract_text(self, page_content):
        try:
            text = page_content.get_text()
            if text != "":
                self.text_dict['page'][self.page_num].append(
                    text.strip())
        except:
            pass

    def _extract_image(self, page_content):
        images_list = page_content.get_images()
        for img_info in images_list:
            # Extract image
            image = self.pdf.extract_image(img_info[0])
            ext = image['ext']
            size = int(image['height']), int(image['width'])
            if ext in ['png', 'jpeg', 'jpg'] and size[0] >= 200 and size[1] >= 200:
                self.img_num += 1
                num = str(self.img_num).zfill(3)
                save_path = f"{self.fileutil.res_dir}{self.name}/image_{num}.{ext}"

                # Store image bytes
                image_blob = image['image']

                # Save image
                with open(save_path, 'wb') as file:
                    file.write(image_blob)
                self.image_blob.append(image_blob)
                self.img_dict[self.img_num].append(save_path)



def extract_info(file_name):
    _, ext = os.path.splitext(file_name)
    ext = ext[1:].lower()

    if ext == 'pptx':
        extractor = PPTInfoExtract()
    elif ext == 'docx':
        extractor = DOCXInfoExtract()
    elif ext == 'pdf':
        extractor = PDFInfoExtract()
    else:
        raise ValueError("Unsupported file format.")

    return extractor.run(file_name)


if __name__ == "__main__":
    file_name = "example.pptx"
    result = extract_info(file_name)
    print(result)
